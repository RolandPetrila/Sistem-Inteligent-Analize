"""
PPTX Generator — python-pptx.
Genereaza prezentare PowerPoint profesionala (7 slide-uri) din verified_data.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


ACCENT = RGBColor(99, 102, 241)
DARK_BG = RGBColor(26, 26, 46)
LIGHT_TEXT = RGBColor(226, 232, 240)
GRAY_TEXT = RGBColor(148, 163, 184)
GREEN = RGBColor(34, 197, 94)
YELLOW = RGBColor(234, 179, 8)
RED = RGBColor(239, 68, 68)
WHITE = RGBColor(255, 255, 255)


def _risk_color(score):
    if isinstance(score, str):
        return {"Verde": GREEN, "Galben": YELLOW, "Rosu": RED}.get(score, GRAY_TEXT)
    if isinstance(score, (int, float)):
        if score >= 70: return GREEN
        if score >= 40: return YELLOW
        return RED
    return GRAY_TEXT


def _add_text(slide, left, top, width, height, text, font_size=14, color=LIGHT_TEXT, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def _set_slide_bg(slide, color=DARK_BG):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def generate_pptx(report_sections: dict, meta: dict, verified_data: dict, output_path: str):
    """Genereaza PPTX din verified_data + meta."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # D13 fix: Guard against None risk_score
    risk_score = verified_data.get("risk_score") or {}
    company = verified_data.get("company") or {}

    # --- Slide 1: Cover ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide1)
    _add_text(slide1, 1, 1, 11, 1, "Roland Intelligence System", 16, GRAY_TEXT, alignment=PP_ALIGN.CENTER)
    _add_text(slide1, 1, 2.2, 11, 1.2, meta.get("title", "Raport"), 36, ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    _add_text(slide1, 1, 3.5, 11, 0.8, meta.get("company_name", ""), 24, LIGHT_TEXT, alignment=PP_ALIGN.CENTER)
    score_text = f"Scor Risc: {risk_score.get('score', 'N/A')}"
    numeric = risk_score.get("numeric_score")
    if numeric is not None:
        score_text += f" ({numeric}/100)"
    _add_text(slide1, 1, 4.8, 11, 0.6, score_text, 20, _risk_color(risk_score.get("score", "")), bold=True, alignment=PP_ALIGN.CENTER)
    _add_text(slide1, 1, 6, 11, 0.5, f"Nivel {meta.get('report_level', 'N/A')} | {meta.get('generated_at', '')} | {meta.get('sources_count', 0)} surse", 12, GRAY_TEXT, alignment=PP_ALIGN.CENTER)

    # --- Slide 2: Profil Firma ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide2)
    _add_text(slide2, 0.5, 0.3, 12, 0.6, "Profil Firma", 28, ACCENT, bold=True)

    profile_fields = [
        ("Denumire", "denumire"), ("CUI", "cui"), ("Adresa", "adresa"),
        ("Nr. Reg. Com.", "nr_reg_com"), ("Stare", "stare_inregistrare"),
        ("Platitor TVA", "platitor_tva"), ("Inactiv", "inactiv"),
    ]
    y = 1.2
    for label, key in profile_fields:
        field = company.get(key, {})
        val = field.get("value", "N/A") if isinstance(field, dict) else "N/A"
        trust = field.get("trust", "") if isinstance(field, dict) else ""
        display = f"{label}: {val}"
        if trust:
            display += f"  [{trust}]"
        _add_text(slide2, 0.8, y, 11, 0.4, display, 14, LIGHT_TEXT)
        y += 0.45

    # --- Slide 3: Date Financiare ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide3)
    _add_text(slide3, 0.5, 0.3, 12, 0.6, "Date Financiare", 28, ACCENT, bold=True)

    financial = verified_data.get("financial", {})
    fin_fields = [
        ("Cifra de Afaceri", "cifra_afaceri"),
        ("Profit Net", "profit_net"),
        ("Nr. Angajati", "numar_angajati"),
        ("Capitaluri Proprii", "capitaluri_proprii"),
        ("Platitor TVA", "platitor_tva"),
        ("Curs EUR/RON", "eur_ron_rate"),
    ]
    y = 1.2
    for label, key in fin_fields:
        field = financial.get(key, {})
        val = field.get("value", "N/A") if isinstance(field, dict) else "N/A"
        note = field.get("note", "") if isinstance(field, dict) else ""
        if isinstance(val, (int, float)) and abs(val) > 1000:
            display = f"{label}: {val:,.0f} RON"
        else:
            display = f"{label}: {val}"
        if note:
            display += f"  ({note})"
        _add_text(slide3, 0.8, y, 11, 0.4, display, 14, LIGHT_TEXT)
        y += 0.45

    # B17: Add financial trend data if available
    trend = financial.get("trend_financiar", {})
    trend_val = trend.get("value") if isinstance(trend, dict) else None
    if isinstance(trend_val, dict):
        y += 0.3
        _add_text(slide3, 0.8, y, 11, 0.4, "Trend Multi-An:", 14, ACCENT, bold=True)
        y += 0.45
        for metric_key in ["cifra_afaceri_neta", "profit_net"]:
            metric_data = trend_val.get(metric_key, {})
            if isinstance(metric_data, dict) and metric_data.get("values"):
                growth = metric_data.get("growth_percent")
                direction = metric_data.get("direction", "")
                name = metric_data.get("name", metric_key)
                color = GREEN if direction == "crestere" else RED if direction == "scadere" else LIGHT_TEXT
                growth_str = f" ({'+' if growth and growth > 0 else ''}{growth}%)" if growth is not None else ""
                _add_text(slide3, 0.8, y, 11, 0.4, f"{name}: {direction}{growth_str}", 12, color)
                y += 0.4

    # --- Slide 4: Scor Risc ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide4)
    _add_text(slide4, 0.5, 0.3, 12, 0.6, "Evaluare Risc", 28, ACCENT, bold=True)

    score_display = f"{risk_score.get('score', 'N/A')}"
    if numeric is not None:
        score_display += f"  —  {numeric}/100"
    _add_text(slide4, 1, 1.2, 11, 0.8, score_display, 40, _risk_color(risk_score.get("score")), bold=True, alignment=PP_ALIGN.CENTER)
    _add_text(slide4, 1, 2.2, 11, 0.5, risk_score.get("recommendation", ""), 14, GRAY_TEXT, alignment=PP_ALIGN.CENTER)

    # Dimensiuni
    dimensions = risk_score.get("dimensions", {})
    if dimensions:
        y = 3.2
        for dim_name, dim_data in dimensions.items():
            sc = dim_data.get("score", 0)
            w = dim_data.get("weight", 0)
            _add_text(slide4, 1, y, 3, 0.35, f"{dim_name.capitalize()} ({w}%)", 12, LIGHT_TEXT)
            _add_text(slide4, 4.5, y, 1.5, 0.35, f"{sc}/100", 12, _risk_color(sc), bold=True)
            y += 0.4

    # Factori
    factors = risk_score.get("factors", [])
    if factors:
        y_f = 3.2
        _add_text(slide4, 7, y_f - 0.4, 5, 0.35, "Factori de Risc:", 13, ACCENT, bold=True)
        for factor, severity in factors[:6]:
            sev_color = {"HIGH": RED, "MEDIUM": YELLOW, "LOW": RGBColor(96, 165, 250), "POSITIVE": RGBColor(74, 222, 128)}.get(severity, GRAY_TEXT)
            _add_text(slide4, 7, y_f, 5.5, 0.35, f"[{severity}] {factor}", 11, sev_color)
            y_f += 0.35

    # --- Slide 5: Anomalii ---
    anomalies = verified_data.get("anomalies", [])
    if anomalies:
        slide5 = prs.slides.add_slide(prs.slide_layouts[6])
        _set_slide_bg(slide5)
        _add_text(slide5, 0.5, 0.3, 12, 0.6, "Alerte si Anomalii", 28, ACCENT, bold=True)
        y = 1.2
        for anomaly in anomalies[:8]:
            level = anomaly.get("level", "INFO")
            icon = {"SUSPECT": "!!!", "ATENTIE": "!!", "INFO": "i"}.get(level, "")
            color = {"SUSPECT": RED, "ATENTIE": YELLOW, "INFO": GRAY_TEXT}.get(level, GRAY_TEXT)
            _add_text(slide5, 0.8, y, 11, 0.35, f"[{level}] {anomaly.get('rule', '')}", 13, color, bold=True)
            _add_text(slide5, 0.8, y + 0.35, 11, 0.4, anomaly.get("detail", ""), 11, LIGHT_TEXT)
            y += 0.85

    # --- Slide 6: Sectiuni raport (rezumat) ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide6)
    _add_text(slide6, 0.5, 0.3, 12, 0.6, "Concluzii", 28, ACCENT, bold=True)
    y = 1.2
    for key, section in list(report_sections.items())[:4]:
        title = section.get("title", key)
        content = section.get("content", "")
        # Primul paragraf ca rezumat
        first_para = ""
        for line in content.split("\n"):
            line = line.strip()
            if line and not line.startswith("**") and not line.startswith("##") and not line.startswith("-"):
                first_para = line[:200]
                break
        _add_text(slide6, 0.8, y, 11, 0.35, title, 14, ACCENT, bold=True)
        if first_para:
            _add_text(slide6, 0.8, y + 0.4, 11, 0.5, first_para, 11, LIGHT_TEXT)
        y += 1.1

    # --- Slide 7: Surse + Disclaimer ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide7)
    _add_text(slide7, 0.5, 0.3, 12, 0.6, "Surse Utilizate", 28, ACCENT, bold=True)

    sources = meta.get("sources", [])
    y = 1.2
    for src in sources[:10]:
        lvl = src.get("level", "?")
        name = src.get("name", "N/A")
        status = src.get("status", "?")
        color = GREEN if status == "OK" else RED
        _add_text(slide7, 0.8, y, 11, 0.3, f"[Nivel {lvl}] {name} — {status}", 11, color)
        y += 0.3

    _add_text(slide7, 0.5, 5.8, 12, 0.8,
        "Disclaimer: Acest raport a fost generat automat folosind exclusiv date publice. "
        "Nu inlocuieste consultanta profesionala.", 9, GRAY_TEXT, alignment=PP_ALIGN.CENTER)
    _add_text(slide7, 0.5, 6.8, 12, 0.4, "Roland Intelligence System v1.1", 10, GRAY_TEXT, alignment=PP_ALIGN.CENTER)

    prs.save(output_path)
