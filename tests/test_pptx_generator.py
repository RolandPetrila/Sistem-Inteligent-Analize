"""B3 (2026-07-16): pptx_generator computed a per-severity `icon` on the anomalies
slide (Slide 5) but never rendered it (ruff F841) — the icon glyph ("!!!"/"!!"/"i")
was thrown away. Decision: RENDER it (prepended to the existing colored/bold
severity label line) rather than delete it — it is a real, cheap visual severity
marker on a slide that is otherwise plain [LEVEL] text, consistent with the
existing style (color + bold already encode severity; the icon reinforces it for
a slide deck skimmed quickly). Verified with a content test via python-pptx (not
just "does not raise")."""

import os
import tempfile


def _verified_data_with_anomalies() -> dict:
    return {
        "company": {"denumire": {"value": "Test SRL"}, "cui": {"value": "26313362"}},
        "risk_score": {"score": "Rosu", "numeric_score": 20, "dimensions": {}, "factors": []},
        "financial": {},
        "anomalies": [
            {"level": "SUSPECT", "rule": "Zombie firma", "detail": "CA=0 si angajati=0"},
            {"level": "ATENTIE", "rule": "Scadere CA", "detail": "CA scazuta cu 40%"},
            {"level": "INFO", "rule": "Nota informativa", "detail": "Fara impact pe scor"},
        ],
    }


def _meta() -> dict:
    return {
        "title": "Raport Test",
        "company_name": "Test SRL",
        "report_level": 2,
        "generated_at": "16.07.2026 10:00",
        "sources_count": 2,
        "sources": [{"name": "ANAF", "level": 1, "status": "OK"}],
    }


class TestPptxAnomalyIcons:
    def test_severity_icons_are_rendered_on_anomaly_slide(self):
        """DOVADA DE NON-VACUITATE: pe codul vechi, `icon` era calculat si niciodata
        concatenat in textul randat -> acest test pica pe codul vechi pentru ca
        textul slide-ului nu contine glifurile de severitate ("!!!"/"!!").
        Verificat cu `git stash` pe codul pre-fix: FAIL (vezi raportul agentului)."""
        from pptx import Presentation

        from backend.reports.pptx_generator import generate_pptx

        report_sections = {
            "executive_summary": {"title": "Rezumat", "content": "Firma analizata."}
        }
        verified_data = _verified_data_with_anomalies()
        meta = _meta()

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            path = f.name
        try:
            generate_pptx(report_sections, meta, verified_data, path)
            assert os.path.exists(path)
            assert os.path.getsize(path) > 0

            prs = Presentation(path)
            # Slide 5 (index 4): "Alerte si Anomalii" — only present because
            # verified_data["anomalies"] is non-empty.
            anomaly_slide = None
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame and "Alerte si Anomalii" in shape.text_frame.text:
                        anomaly_slide = slide
                        break
                if anomaly_slide:
                    break
            assert anomaly_slide is not None, "Slide-ul de anomalii nu a fost gasit"

            slide_text = "\n".join(
                shape.text_frame.text for shape in anomaly_slide.shapes if shape.has_text_frame
            )

            assert "!!!" in slide_text  # SUSPECT icon
            assert "!!" in slide_text  # ATENTIE icon (substring of "!!!", checked separately below)
            assert "[SUSPECT] Zombie firma" in slide_text
            assert "[ATENTIE] Scadere CA" in slide_text
            assert "[INFO] Nota informativa" in slide_text

            # Distinguish ATENTIE ("!!") from SUSPECT ("!!!") precisely, line by line.
            lines = [line for line in slide_text.split("\n") if "[SUSPECT]" in line or "[ATENTIE]" in line or "[INFO]" in line]
            suspect_line = next(line for line in lines if "[SUSPECT]" in line)
            atentie_line = next(line for line in lines if "[ATENTIE]" in line)
            info_line = next(line for line in lines if "[INFO]" in line)
            assert suspect_line.strip().startswith("!!!")
            assert atentie_line.strip().startswith("!!") and not atentie_line.strip().startswith("!!!")
            assert info_line.strip().startswith("i ")
        finally:
            if os.path.exists(path):
                os.remove(path)

    def test_no_anomalies_skips_slide_without_raising(self):
        """Edge case: verified_data['anomalies'] gol -> slide-ul de anomalii nu
        trebuie generat, iar functia nu trebuie sa arunce."""
        from pptx import Presentation

        from backend.reports.pptx_generator import generate_pptx

        report_sections = {"executive_summary": {"title": "Rezumat", "content": "Firma analizata."}}
        verified_data = {
            "company": {"denumire": {"value": "Test SRL"}},
            "risk_score": {"score": "Verde", "numeric_score": 85, "dimensions": {}, "factors": []},
            "financial": {},
            "anomalies": [],
        }
        meta = _meta()

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            path = f.name
        try:
            generate_pptx(report_sections, meta, verified_data, path)
            assert os.path.exists(path)

            prs = Presentation(path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        assert "Alerte si Anomalii" not in shape.text_frame.text
        finally:
            if os.path.exists(path):
                os.remove(path)
