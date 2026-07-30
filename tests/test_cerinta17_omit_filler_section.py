"""CERINTA #17 (P6): omite din raport sectiunile "date insuficiente" (filler marcat de sinteza).

Design = mark-and-skip la randare (NU don't-emit): sinteza pune `INSUFFICIENT_DATA_MARKER=True` pe
fillerul determinist la UN SINGUR punct de emisie; cheia ramane in `report_sections` (deci in DB /
regenerare / chat), iar cei 4 randere o OMIT prin `visible_sections`.

Non-vacuitate (dovada mecanica in JURNAL, swap git-show-HEAD): pe HEAD randerele itereaza
`report_sections` brut -> fillerul marcat SE RANDEAZA -> testele de ABSENTA pica. Producatorul pe
HEAD nu pune marker -> testul de marker pica.

Invariant #4 (never-empty) = test de GARDA (trece si pe cod vechi): protejeaza contra unei
implementari naive "sari mereu", care ar produce un raport gol cand toate sectiunile-s filler.
"""
import os
import tempfile

import pytest

from backend.reports.docx_generator import generate_docx
from backend.reports.html_generator import generate_html
from backend.reports.pdf_generator import generate_pdf
from backend.reports.pptx_generator import generate_pptx
from backend.reports.section_visibility import (
    INSUFFICIENT_DATA_MARKER,
    is_filler_section,
    visible_sections,
)

# Santinela: un singur token ASCII care NU apare nicaieri altundeva in raport -> absenta neambigua
# (evita fals-pozitive din TOC/proza care ar contine "Competi..."). Truc reutilizat din #16 D (ZZZ).
SENTINEL_TITLE = "ZZZCOMPETITIONSENTINEL"
EXEC_TITLE = "ZZZEXECSENTINEL"
REC_TITLE = "ZZZRECOMMENDSENTINEL"  # a 5-a sectiune reala -> proba inv.#3 (filter-inainte-de-slice PPTX)


def _filler_competition() -> dict:
    """Fillerul "date insuficiente" asa cum il emite sinteza (cu marker)."""
    return {
        "title": SENTINEL_TITLE,
        "content": "Sectiunea nu a putut fi generata din cauza datelor insuficiente disponibile.",
        "word_count": 0,
        INSUFFICIENT_DATA_MARKER: True,
    }


def _real_competition() -> dict:
    """ACELASI titlu santinela, DAR fara marker (continut real) -> trebuie sa ramana randata."""
    return {
        "title": SENTINEL_TITLE,
        "content": "Concurentii principali sunt firmele Alfa si Beta cu cote de piata relevante.",
        "word_count": 42,
    }


def _real_exec() -> dict:
    return {"title": EXEC_TITLE, "content": "Firma este stabila financiar si operational."}


def _real_recommendations() -> dict:
    return {"title": REC_TITLE, "content": "Se recomanda diversificarea portofoliului de clienti."}


def _plain(title: str) -> dict:
    return {"title": title, "content": f"Continut real pentru {title}."}


def _meta() -> dict:
    return {
        "title": "Raport Test",
        "company_name": "Test SRL",
        "report_level": 3,
        "generated_at": "30.07.2026 10:00",
        "risk_score": "Verde",
        "numeric_score": 80,
        "risk_recommendation": "",
        "sources_count": 1,
        "sources": [],
        "report_number": "RIS-2026-0001",
    }


def _vd() -> dict:
    return {"risk_score": {}, "completeness": {}}


def _pdf_sanitize_mirror(obj):
    """Oglinda EXACTA a closure-ului `_pdf_sanitize` din generator.py:86-93 = calea reala PDF.
    PDF e SINGURUL renderer care primeste un dict TRANSFORMAT (sanitizat latin-1); un test care ar
    apela `generate_pdf(report_sections)` direct ar rata un bug care sterge markerul la sanitizare
    (clasa-semnatura a proiectului). Aici dovedim ca markerul (cheie ASCII, valoare bool) SUPRAVIETUIESTE.
    """
    from backend.reports.pdf_generator import _sanitize as pdf_sanitize_fn
    if isinstance(obj, str):
        return pdf_sanitize_fn(obj)
    if isinstance(obj, dict):
        return {_pdf_sanitize_mirror(k): _pdf_sanitize_mirror(v) for k, v in obj.items()}
    if isinstance(obj, list):
        return [_pdf_sanitize_mirror(i) for i in obj]
    return obj


def _pdf_text(path: str) -> str:
    import pdfplumber
    with pdfplumber.open(path) as pdf:
        return "\n".join(page.extract_text() or "" for page in pdf.pages)


def _docx_text(path: str) -> str:
    from docx import Document
    return "\n".join(p.text for p in Document(path).paragraphs)


def _pptx_text(path: str) -> str:
    from pptx import Presentation
    out = []
    for slide in Presentation(path).slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                out.append(shape.text_frame.text)
    return "\n".join(out)


# ---------------------------------------------------------------------------
# Unit: visible_sections / is_filler_section
# ---------------------------------------------------------------------------

class TestVisibilityHelpers:
    def test_marked_filler_is_filler(self):
        assert is_filler_section(_filler_competition()) is True

    def test_unmarked_section_not_filler(self):
        assert is_filler_section(_real_competition()) is False

    def test_word_count_zero_without_marker_not_filler(self):
        # INVARIANT #1: cheiaza pe MARKER, NU pe word_count==0. O sectiune degradata de deadline
        # (word_count 0, FARA marker) trebuie PASTRATA — nu e "date insuficiente".
        degraded = {"title": "X", "content": "text degradat", "word_count": 0}
        assert is_filler_section(degraded) is False

    def test_string_value_not_filler(self):
        # Garda portanta: orchestrator.py pune o valoare STRING pe calea de eroare -> nu trebuie crash.
        assert is_filler_section("Eroare generare raport: ...") is False

    def test_none_not_filler(self):
        assert is_filler_section(None) is False

    def test_visible_drops_filler_keeps_order(self):
        sections = {
            "executive_summary": _real_exec(),
            "competition": _filler_competition(),
            "recommendations": {"title": "Recomandari", "content": "..."},
        }
        vis = visible_sections(sections)
        assert list(vis.keys()) == ["executive_summary", "recommendations"]

    def test_visible_never_empty_when_all_filler(self):
        # INVARIANT #4: toate filler -> pastreaza tot (fillerul e raspunsul onest, nu body gol).
        sections = {"competition": _filler_competition()}
        assert visible_sections(sections) == sections

    def test_visible_non_dict_passthrough(self):
        assert visible_sections("not a dict") == "not a dict"


# ---------------------------------------------------------------------------
# Non-vac: randerele OMIT fillerul marcat din raport (HTML / PDF / DOCX)
# ---------------------------------------------------------------------------

class TestFillerOmittedFromRenderers:
    def test_html_omits_marked_filler(self):
        sections = {"executive_summary": _real_exec(), "competition": _filler_competition()}
        with tempfile.NamedTemporaryFile(suffix=".html", delete=False) as f:
            path = f.name
        try:
            generate_html(sections, _meta(), _vd(), path)
            with open(path, encoding="utf-8") as fh:
                content = fh.read()
            assert SENTINEL_TITLE not in content        # sectiunea filler nu se randeaza
            assert EXEC_TITLE in content                 # sectiunea reala ramane
            assert 'id="competition"' not in content     # INVARIANT #2: zero ancora moarta
            assert 'href="#competition"' not in content  # INVARIANT #2: zero link nav mort
        finally:
            if os.path.exists(path):
                os.remove(path)

    def test_html_keeps_real_section(self):
        sections = {"executive_summary": _real_exec(), "competition": _real_competition()}
        with tempfile.NamedTemporaryFile(suffix=".html", delete=False) as f:
            path = f.name
        try:
            generate_html(sections, _meta(), _vd(), path)
            with open(path, encoding="utf-8") as fh:
                content = fh.read()
            assert SENTINEL_TITLE in content             # fara marker -> se randeaza normal
            assert 'href="#competition"' in content
        finally:
            if os.path.exists(path):
                os.remove(path)

    def test_pdf_omits_marked_filler_through_real_sanitize(self):
        sections = {"executive_summary": _real_exec(), "competition": _filler_competition()}
        sanitized = _pdf_sanitize_mirror(sections)
        # Dovada: markerul supravietuieste sanitizarii latin-1 (calea reala PDF).
        assert sanitized["competition"].get(INSUFFICIENT_DATA_MARKER) is True
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
            path = f.name
        try:
            generate_pdf(sanitized, _pdf_sanitize_mirror(_meta()), path, verified_data=_vd())
            text = _pdf_text(path)
            assert SENTINEL_TITLE not in text
            assert EXEC_TITLE in text
        finally:
            if os.path.exists(path):
                os.remove(path)

    def test_pdf_keeps_real_section_through_real_sanitize(self):
        sections = {"executive_summary": _real_exec(), "competition": _real_competition()}
        sanitized = _pdf_sanitize_mirror(sections)
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
            path = f.name
        try:
            generate_pdf(sanitized, _pdf_sanitize_mirror(_meta()), path, verified_data=_vd())
            text = _pdf_text(path)
            assert SENTINEL_TITLE in text
        finally:
            if os.path.exists(path):
                os.remove(path)

    def test_docx_omits_marked_filler(self):
        sections = {"executive_summary": _real_exec(), "competition": _filler_competition()}
        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            path = f.name
        try:
            generate_docx(sections, _meta(), path, verified_data=_vd())
            text = _docx_text(path)
            assert SENTINEL_TITLE not in text
            assert EXEC_TITLE in text
        finally:
            if os.path.exists(path):
                os.remove(path)

    def test_docx_keeps_real_section(self):
        sections = {"executive_summary": _real_exec(), "competition": _real_competition()}
        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            path = f.name
        try:
            generate_docx(sections, _meta(), path, verified_data=_vd())
            text = _docx_text(path)
            assert SENTINEL_TITLE in text
        finally:
            if os.path.exists(path):
                os.remove(path)


# ---------------------------------------------------------------------------
# Invariant #4 (garda): daca fillerul e SINGURA sectiune -> ramane randat (never-empty body)
# ---------------------------------------------------------------------------

class TestNeverEmptyBody:
    def test_html_keeps_lone_filler(self):
        sections = {"competition": _filler_competition()}
        with tempfile.NamedTemporaryFile(suffix=".html", delete=False) as f:
            path = f.name
        try:
            generate_html(sections, _meta(), _vd(), path)
            with open(path, encoding="utf-8") as fh:
                content = fh.read()
            assert SENTINEL_TITLE in content  # nu se omite -> body-ul nu ramane gol
        finally:
            if os.path.exists(path):
                os.remove(path)


# ---------------------------------------------------------------------------
# Invariant #3 (PPTX): filtreaza INAINTE de slice [:4] (decizie declarata).
# Pune fillerul in primele 4 -> nu e "sters" doar de slice, ci de filtrare.
# ---------------------------------------------------------------------------

class TestPptxFilterBeforeSlice:
    def test_pptx_filters_filler_then_slices(self):
        # Ordine: exec(1), competition-filler(2), financial(3), risk(4), recommendations(5).
        # NEW: visible_sections scoate competition -> [:4] = exec/financial/risk/recommendations
        #      -> competition ABSENT + recommendations (a 5-a) PROMOVATA in slide-ul Concluzii.
        # HEAD: [:4] brut = exec/competition/financial/risk -> competition PREZENT (slot 2) +
        #       recommendations ABSENT (taiat de slice) -> AMBELE aserturi pica pe HEAD.
        sections = {
            "executive_summary": _real_exec(),
            "competition": _filler_competition(),
            "financial_analysis": _plain("ZZZFINSENTINEL"),
            "risk_assessment": _plain("ZZZRISKSENTINEL"),
            "recommendations": _real_recommendations(),
        }
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            path = f.name
        try:
            generate_pptx(sections, _meta(), _vd(), path)
            text = _pptx_text(path)
            assert SENTINEL_TITLE not in text   # fillerul filtrat inainte de slice
            assert REC_TITLE in text            # a 5-a sectiune reala promovata in Concluzii
        finally:
            if os.path.exists(path):
                os.remove(path)


# ---------------------------------------------------------------------------
# Producator: sinteza marcheaza fillerul la punctul unic de emisie
# ---------------------------------------------------------------------------

class TestSynthesisMarksFiller:
    @pytest.mark.asyncio
    async def test_generate_section_marks_insufficient_competition(self):
        import backend.agents.agent_synthesis as m

        agent = m.SynthesisAgent()
        # competition FARA web_presence.competitors -> _has_sufficient_data False -> filler.
        # Returneaza inainte de orice apel de provider AI (ramura ER2), deci nu mockam nimic.
        section = {"key": "competition", "title": SENTINEL_TITLE, "word_count": 400}
        result = await agent.generate_section(section, {"completeness": {"score": 90}})
        assert result.get(INSUFFICIENT_DATA_MARKER) is True
        assert result["word_count"] == 0
